"""
Generates docs/Project_Presentation.pptx from the project's actual training results.
Run from docs/: python generate_pptx.py
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
METRICS_PATH = os.path.join(ROOT, "saved_models", "metrics.json")

PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0A, 0x0A, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    return slide


def add_title(slide, text, size=36, top=0.5, color=WHITE):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12), Inches(1))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return box


def add_bullets(slide, items, top=1.6, size=20, width=11.5, color=WHITE):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(width), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return box


def accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()


# Slide 1: Title
s = add_slide()
accent_bar(s)
add_title(s, "Credit Card Fraud Detection AI System", size=40, top=2.6)
sub = s.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(12), Inches(1))
sub.text_frame.text = "Real-Time Fraud Detection Platform — ML + Business Rule Engine"
sub.text_frame.paragraphs[0].font.size = Pt(22)
sub.text_frame.paragraphs[0].font.color.rgb = GRAY
sub2 = s.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.6))
sub2.text_frame.text = "Internship Submission — Full-Stack ML Engineering Project"
sub2.text_frame.paragraphs[0].font.size = Pt(16)
sub2.text_frame.paragraphs[0].font.color.rgb = GRAY

# Slide 2: Problem Statement
s = add_slide(); accent_bar(s)
add_title(s, "Problem Statement")
add_bullets(s, [
    "Credit card fraud costs the global financial industry tens of billions of dollars annually.",
    "Traditional rule-only systems miss novel fraud patterns and are hard to tune.",
    "Pure ML systems are black boxes — hard to audit or explain to compliance/regulators.",
    "Goal: combine ML risk scoring with an explicit, auditable rule engine, in a real-time system a fraud operations team could actually use.",
])

# Slide 3: Architecture
s = add_slide(); accent_bar(s)
add_title(s, "System Architecture")
add_bullets(s, [
    "Frontend: React + Vite + Tailwind + Framer Motion (7 pages, dark/light theme)",
    "Backend: FastAPI — JWT auth, rate limiting, structured logging, 18 REST endpoints",
    "ML Layer: feature engineering → trained model → SHAP explainability",
    "Decision Engine: combines ML probability with explicit business rules",
    "Database: SQLite by default, PostgreSQL via DATABASE_URL — 7 tables",
], size=19)

# Slide 4: Data
s = add_slide(); accent_bar(s)
add_title(s, "Data Strategy")
add_bullets(s, [
    "Synthetic business-schema dataset — powers the interactive Predict form (human-meaningful fields)",
    "Real Kaggle Credit Card Fraud dataset (anonymized V1-V28 PCA features) — separate benchmark pipeline",
    "Feature engineering: amount_ratio, velocity scores, geo-distance, device/location change flags, merchant frequency",
    "SMOTE oversampling applied to the training split only — no leakage into evaluation",
])

# Slide 5: Model Comparison (real data)
s = add_slide(); accent_bar(s)
add_title(s, "Model Comparison (Actual Training Run)")
if os.path.exists(METRICS_PATH):
    with open(METRICS_PATH) as f:
        metrics = json.load(f)
    rows = metrics["results"]
    best = metrics["best_model"]

    table_shape = s.shapes.add_table(len(rows) + 1, 6, Inches(0.7), Inches(1.6), Inches(12), Inches(0.5 * (len(rows) + 1)))
    table = table_shape.table
    headers = ["Model", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE

    for i, (name, m) in enumerate(rows.items(), start=1):
        vals = [name, f"{m['accuracy']:.3f}", f"{m['precision']:.3f}", f"{m['recall']:.3f}", f"{m['f1']:.3f}", f"{m['roc_auc']:.3f}"]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = v
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE if name != best else RGBColor(0x22, 0xC5, 0x5E)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x24)

    note = s.shapes.add_textbox(Inches(0.7), Inches(1.6 + 0.5 * (len(rows) + 1) + 0.2), Inches(12), Inches(0.5))
    note.text_frame.text = f"Best model selected: {best} (highlighted in green)"
    note.text_frame.paragraphs[0].font.size = Pt(14)
    note.text_frame.paragraphs[0].font.color.rgb = GRAY

# Slide 6: Decision Engine
s = add_slide(); accent_bar(s)
add_title(s, "ML + Rule Engine Decision Flow")
add_bullets(s, [
    "ML model produces a probability from learned patterns",
    "Rule engine evaluates explicit business rules (velocity bursts, impossible travel, high-risk merchant, VPN, failed OTP, new device+location+high amount)",
    "Combined score = 1 - (1 - ml_prob) × (1 - rule_score)",
    "Risk bucketed: Low / Medium / High / Critical",
    "Recommended action: Approve / Send OTP / Manual Review / Decline / Freeze Account",
    "Plain-English explanation generated for every decision",
])

# Slide 7: Frontend
s = add_slide(); accent_bar(s)
add_title(s, "Frontend — Fraud Operations Center")
add_bullets(s, [
    "Dashboard: 7 KPIs, fraud trend/hour/device/age/location/category charts",
    "Predict: full banking transaction simulator with animated speedometer risk meter",
    "Alert Center: live queue with Approve/Block/Review/Freeze actions",
    "History: search, filters, detail modal, CSV + PDF export",
    "Customer Profiles: spending history, merchant breakdown, fraud history",
    "Explainability: SHAP visualizations + plain-English reasoning",
])

# Slide 8: Testing
s = add_slide(); accent_bar(s)
add_title(s, "Testing & Validation")
add_bullets(s, [
    "30 realistic scenarios tested end-to-end against the live API",
    "Routine purchases (grocery, fuel, subscriptions) → consistently Low risk, auto-approved",
    "Compounding risk signals (new device + location + high amount + failed OTP) → Critical, frozen",
    "Single weak signals alone do not trigger auto-decline — mirrors real fraud-ops false-positive tolerance",
    "See docs/TESTING_REPORT.md for full results",
])

# Slide 9: Security & Ops
s = add_slide(); accent_bar(s)
add_title(s, "Security & Production Readiness")
add_bullets(s, [
    "JWT authentication, bcrypt password hashing",
    "Rate limiting (120 req/min per IP), input validation via Pydantic",
    "Structured logging, global exception handler with clean error responses",
    "Configurable CORS, environment-variable based secrets",
    "SQLite for zero-setup local dev; PostgreSQL for production via DATABASE_URL",
    "Dockerized: docker-compose up brings up Postgres + backend + frontend",
])

# Slide 10: Conclusion
s = add_slide(); accent_bar(s)
add_title(s, "Conclusion & Future Work")
add_bullets(s, [
    "Built a complete, verified, end-to-end real-time fraud detection platform",
    "Combines interpretable ML with an auditable business rule engine",
    "Future: streaming ingestion, MLflow model registry, RBAC, Redis-backed rate limiting at scale",
    "Thank you — questions?",
])

out_path = os.path.join(ROOT, "docs", "Project_Presentation.pptx")
prs.save(out_path)
print("Saved", out_path)
